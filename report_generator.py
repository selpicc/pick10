# -*- coding: utf-8 -*-
"""브랜드별 셀픽 Media Proposal(분석 리포트) PPTX 생성기
────────────────────────────────────────────────────────────────
대시보드 셀러 디테일의 '📊 분석 리포트 만들기' 버튼이 이 모듈을 호출한다.
자동화(스케줄러)가 아니라 사람이 버튼을 누른 그 브랜드에 대해서만 1건 생성한다.

만드는 것: 표지 없이 3페이지짜리 '기획서'형 편집 가능 PPTX (폰트 Pretendard).
  1) 컨셉·전략 = 슬로건 + 빅 아이디어(컨셉) + 시장 인사이트 + 핵심 타겟
  2) 시장·니즈 = 네이버 실시간 검색(진짜 수치) + 실제 언급 키워드 + 소구 키워드
  3) 셀픽 실행 = 광고 상품 5종 + 오프/온 실행 아이디어 + 도달 규모
  ※ 준 예시 PDF는 '내용/분석' 참고용이며 레이아웃은 새로 잡는다.
  ※ 외부 시장 데이터는 네이버 검색 실측만 사용(가짜 수치 X). 빅아이디어·컨셉은 창작.

핵심 원칙 — 지어내지 않기 (brand_intro.py와 동일 철학):
  · 셀픽 고정 수치(조리원 230곳·산부인과 55곳 등)는 소개서 값 그대로 '템플릿'.
  · 브랜드/상품 부분(사용 루틴·타겟·키워드)은 그 브랜드 실제 제품정보 기반 AI 생성.
    검증 불가 주장(1위·인증·특허·기관명 등)이 섞이면 자동 폐기.
  · 창의적 마케팅 아이디어는 슬라이드에 '제안'으로 명확히 표기.
  · AI가 실패/빈약해도 절대 멈추지 않는다 → 카테고리 기반 템플릿으로 완성본 생성.

사용:
    from report_generator import make_report
    pptx_bytes, filename = make_report(row)   # row = sellers 테이블 한 행(dict)
"""
import io
import os
import re
import json

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Gemini 호출·홈페이지 읽기는 brand_intro의 것을 재사용 (중복 방지, 같은 안전 철학)
from brand_intro import _fetch_text, API_KEY, MODEL, API_URL
import requests


# ─────────────────────────────────────────────────────────────
# 셀픽 고정 사실 (MEDIA PROPOSAL 2026 — 지어내지 말 것. 이 값만 사용)
# ─────────────────────────────────────────────────────────────
SELPIC_FACTS = {
    "산모DB": "연간 10만명 신생아·부모 DB 확보 (월 8,000~10,000명 신규가입)",
    "누적회원": "최근 3년 누적 회원(0~3세) 41만명 · 누적 다운로드 100만+",
    "설치점": "전국 460여 대 키오스크 설치 (조리원·산부인과·마트·병원 등)",
    "제휴": "제휴 산후조리원·산부인과 전국 285곳",
    "조리원": "산후조리원 230곳 우선 설치 (시장점유 약 50%)",
    "산부인과": "분만 수 상위 산부인과 55곳 설치",
    "체류": "조리원 체류 14일 매일 접촉 · 산부인과 임신~출산 반복 노출",
    "무료인화": "임산부 1인당 무료 사진 15장 인화 (초음파~아기 사진)",
    "키오스크": "인화 대기 중 100% 노출 · 최소 일 10만 회 노출",
}

# 셀픽 영유아 시장 타겟 특성 (모든 브랜드 공통 — 소개서 기반, 고정)
TARGET_MARKET_LINES = [
    "영유아 시장 = 20대 후반~40대 초반의 여성",
    "엄마가 된다는 것은 삶의 안정기에 접어든 여성 — 좋고 안전한 제품에 아낌이 없다",
    "맘카페 빅마우스이자, 조리원 동기·산후조리원 원장의 추천에 크게 의지한다",
    "본인 경험을 맘카페로 나눠 후배 맘에게 전파 → 그 자체가 바이럴이자 새 시장 창출",
    "온라인에서 탐색·구매하지만, 오프라인의 전문가 사용·실사용 확인이 신뢰의 바탕",
    "3년이면 소비자가 교체되는 짧은 주기 → 지속 마케팅 + 빠른 신규 시장 형성",
]

# 타겟 세그먼트 후보 (예시 덱과 동일 풀)
SEGMENT_POOL = ["산모", "아기", "가족", "청소년", "아토피 환우",
                "새로운 뷰티 트렌드", "청결 강박", "미세먼지"]

# 셀픽 미디어 상품 5종 (2026 가이드 — 고정, 지어내지 말 것)
SELPIC_MEDIA = [
    ("키오스크 DA", "인화 대기 화면 100% 노출 사이니지 광고"),
    ("모바일 DA", "셀픽 앱·모바일 최적 배너 (자체 광고엔진)"),
    ("조리원 MRO·샘플링", "신생아실 사용 물품 결합 오프라인 홍보"),
    ("맘카페 체험형 바이럴", "실제 아기엄마 후기 기반 심화 바이럴"),
    ("타겟 DB LMS", "동의 기반 산모 DB 다이렉트 메시지"),
]

# 검증 불가 주장 — 리포트 문구에 하나라도 있으면 그 항목 폐기
_BANNED_CLAIM = (
    "1위", "일위", "최고", "최초", "유일", "독보", "특허", "인증", "공인",
    "수상", "fda", "식약처", "임상", "의약품", "대웅", "존슨", "공식지정",
    "검증된", "보장", "100%", "무조건 안전",
)


# ─────────────────────────────────────────────────────────────
# 1) 제품 컨텍스트 수집
# ─────────────────────────────────────────────────────────────
def _best_url(row: dict) -> str:
    """홈페이지 텍스트를 읽을 최적 URL. 공식몰 소스가 있으면 우선."""
    src = (row.get("auto_biz_sources") or "")
    m = re.search(r"https?://[^\s,]+", src)
    if m:
        return m.group(0)
    return (row.get("smartstore_url") or "").strip()


def _product_context(row: dict) -> dict:
    brand = (row.get("brand_name") or "").strip()
    category = (row.get("category") or row.get("product_category") or "").strip()
    flagship = (row.get("flagship_product") or "").strip()
    url = _best_url(row)
    site_text = _fetch_text(url) if url else ""
    return {
        "brand": brand,
        "category": category,
        "flagship": flagship,
        "site_text": site_text,
    }


# ─────────────────────────────────────────────────────────────
# 1-2) 외부 시장 데이터 (네이버 실시간 검색) — 진짜 수치만, 실패 시 0/빈값
#   블로그·카페 노출 건수(실제) + 실제 글 제목(사람들이 진짜 쓰는 표현)을 모아
#   AI 컨셉의 '근거'로 넣는다. 지어낸 숫자는 절대 넣지 않는다.
# ─────────────────────────────────────────────────────────────
def _naver_api(kind: str, query: str, display: int = 1) -> dict:
    nid = os.getenv("NAVER_CLIENT_ID", "").strip()
    nsc = os.getenv("NAVER_CLIENT_SECRET", "").strip()
    if not (nid and nsc):
        return {}
    try:
        r = requests.get(
            f"https://openapi.naver.com/v1/search/{kind}.json",
            headers={"X-Naver-Client-Id": nid, "X-Naver-Client-Secret": nsc},
            params={"query": query, "display": display}, timeout=8,
        )
        if r.status_code == 200:
            return r.json()
    except Exception:
        pass
    return {}


def _naver_total(kind: str, q: str) -> int:
    try:
        return int(_naver_api(kind, q, 1).get("total", 0) or 0)
    except Exception:
        return 0


def _naver_titles(kind: str, q: str, n: int) -> list:
    out = []
    for it in (_naver_api(kind, q, n).get("items", []) or []):
        t = re.sub(r"<[^>]+>", "", it.get("title", "") or "")
        t = re.sub(r"&[a-zA-Z#0-9]+;", " ", t)
        t = re.sub(r"\s+", " ", t).strip()
        if t:
            out.append(t)
    return out


# 상품 핵심어 추출 시 버릴 단어(타겟·부가어). 상품 '유형어'만 남긴다.
_PK_STOPWORDS = {
    "신생아", "아기", "유아", "영유아", "베이비", "임산부", "산모", "아이", "키즈",
    "신생", "유아용", "아기용", "아동", "남아", "여아",
    "정품", "공식", "대용량", "세트", "기획", "증정", "사은품", "무료배송", "무배",
    "본품", "리필", "리필용", "선물", "선물세트", "특가", "할인", "구성", "단품", "택",
}


# 상품 '유형어' 사전 — 이 제품이 대체 무엇인지 한 단어로 말해주는 말들.
#   ⭐ 2026-07-20: 예전엔 '남은 토큰 중 뒤 2단어'를 핵심어로 썼는데,
#      스마트스토어 상품명은 유형어가 뒤에 온다는 보장이 없어 자주 빗나갔다.
#      (앙또미뇽 → '올인원 약산성' / 누엔쁘 순환크림 → '화장품 출산선물')
#      그래서 유형어를 먼저 사전으로 찾고, 못 찾을 때만 옛 방식으로 폴백한다.
#   토큰에 '포함'되면 매칭 → '순환크림'은 '크림'으로 걸려 토큰 전체가 핵심어가 된다.
_PRODUCT_TYPES = (
    # 세정·목욕
    "바스앤샴푸", "베이비워시", "바디워시", "클렌징폼", "클렌저", "샴푸", "워시",
    "비누", "입욕제", "버블바스",
    # 보습·스킨케어
    "선크림", "튼살크림", "순환크림", "크림", "로션", "에센스", "세럼", "앰플",
    "미스트", "오일", "밤", "보습제", "수분젤",
    # 구강·위생
    "치약", "칫솔", "가글", "물티슈", "기저귀", "면봉", "손소독",
    # 수유·이유
    "분유", "이유식", "젖병세정제", "젖병", "젖꼭지", "유축기", "수유패드",
    "수유쿠션", "빨대컵", "이유식기", "쪽쪽이",
    # 임산부
    "엽산", "유산균", "비타민", "영양제", "마사지오일",
    # 세탁·생활
    "섬유유연제", "세탁세제", "세제", "살균", "소독",
    # 의류·용품
    "배냇저고리", "우주복", "바디슈트", "턱받이", "속싸개", "겉싸개", "블랭킷",
    "유모차", "카시트", "아기띠", "역류방지쿠션", "베개", "매트",
)
# 긴 유형어부터 검사 (‘바스앤샴푸’가 ‘샴푸’보다 먼저 잡히게)
_PRODUCT_TYPES_SORTED = tuple(sorted(_PRODUCT_TYPES, key=len, reverse=True))


def _clean_tokens(brand: str, flagship: str) -> list:
    """상품명에서 브랜드·타겟어·용량/수량·부가어를 걷어낸 토큰 목록."""
    t = re.sub(r"<[^>]+>", " ", flagship)
    t = re.sub(r"&[a-zA-Z#0-9]+;", " ", t)
    t = re.sub(r"[\[\]（）()\{\}·|/]", " ", t)
    bn = brand.replace(" ", "")
    toks = []
    for tok in t.split():
        tok = tok.strip(",.·")
        if not tok or len(tok) < 2:
            continue
        if brand and (tok in brand or brand in tok):
            continue
        if bn and (tok in bn or bn in tok.replace(" ", "")):
            continue
        if any(ch.isdigit() for ch in tok):     # 280g, 500ml, 1개, 1+1 …
            continue
        if tok in _PK_STOPWORDS:
            continue
        toks.append(tok)
    return toks


def _product_keyword(brand: str, flagship: str) -> str:
    """지저분한 스마트스토어 상품명 → 검색에 쓸 '상품 유형' 핵심어 1개.

    예) '앙또미뇽 아기 바디워시 신생아 유아 키즈 바스앤샴푸 올인원 대용량 약산성 500ml'
        → '바스앤샴푸'   (예전엔 '올인원 약산성' — 제품이 뭔지 안 드러났다)
        '누엔쁘 맘조리 산모 산후 순환크림 임산부 화장품 출산선물'
        → '순환크림'     (예전엔 '화장품 출산선물')

    ⭐ 한 단어만 쓰는 이유: 네이버 검색은 토큰이 늘수록 결과가 급감한다.
       '브랜드 + 유형어' 조합이 "이 브랜드의 이 제품 라인" 반응을 가장 잘 잡는다.
    """
    if not flagship:
        return ""
    toks = _clean_tokens(brand, flagship)
    if not toks:
        return ""

    # 1) 유형어 사전 매칭 — 가장 구체적인(긴 유형어에 걸린) 토큰을 고른다
    best, best_score = "", ()
    for tok in toks:
        for typ in _PRODUCT_TYPES_SORTED:
            if typ in tok:
                score = (len(typ), len(tok))    # 유형어가 길수록·토큰이 길수록 구체적
                if score > best_score:
                    best, best_score = tok, score
                break                           # 이 토큰은 가장 긴 유형어로 이미 매칭됨
    if best:
        return best

    # 2) 사전에 없는 제품(새 유형·잡화) → 옛 방식으로 폴백 (뒤 2단어)
    return " ".join(toks[-2:])


def _market_signals(brand: str, category: str, flagship: str = "") -> dict:
    """이 브랜드 '이 주력상품'에 대한 시장 인식. 전부 실시간 검색(진짜 수치).

    ⭐ 2026-07-20: 기획서는 브랜드가 아니라 '브랜드의 특정 상품' 제안서다.
       그래서 블로그·카페 수치를 브랜드명만이 아니라 '브랜드 + 상품 핵심어'로 센다.
         blog / cafe = "브랜드" + 유형어  (이 제품 라인이 얼마나 회자되나)
       주력상품이 없으면 예전처럼 브랜드 전반으로 폴백한다.

    ⚠ 한때 '시장 전체'(브랜드 뺀 유형어) 칸을 뒀다가 뺐다.
       '아기 바스앤샴푸' 같은 수치는 해당 브랜드 얘기가 아니라서,
       "이 브랜드 주력상품의 시장 인식"을 보려는 목적에 맞지 않았다.
    """
    base = {"blog": 0, "cafe": 0, "prod_kw": "", "brand_only": True, "titles": []}
    if not brand:
        return base

    kw = _product_keyword(brand, flagship)
    if kw and kw not in brand:
        # ── 상품 라인 기준 (기본) ──
        base["prod_kw"] = kw
        base["brand_only"] = False
        pq = f'"{brand}" {kw}'
        base["blog"] = _naver_total("blog", pq)
        base["cafe"] = _naver_total("cafearticle", pq)
        # 글 제목도 '수집 때 저장된 주력상품' 기준으로만 모은다.
        #   브랜드 전반 글을 섞으면 다른 제품 얘기가 표현에 끼어든다
        #   (앙또미뇽 바스앤샴푸 기획서에 콧물흡입기·아토패치 표현이 뽑히던 문제)
        #   주력상품 글이 없으면 표현을 비운다 — 지어내느니 라벨을 '일반'으로 바꾼다.
        titles = (_naver_titles("blog", f"{brand} {kw}", 15)
                  + _naver_titles("cafearticle", f"{brand} {kw}", 12))
    else:
        # ── 주력상품이 없으면 브랜드 전반으로 폴백 ──
        q = f'"{brand}"'
        base["blog"] = _naver_total("blog", q)
        base["cafe"] = _naver_total("cafearticle", q)
        titles = (_naver_titles("blog", brand, 10)
                  + _naver_titles("cafearticle", brand, 10)
                  + _naver_titles("blog", f"{brand} 후기", 6))

    seen, uniq = set(), []
    for t in titles:
        if t not in seen:
            seen.add(t)
            uniq.append(t)
    base["titles"] = uniq[:30]
    return base


# ─────────────────────────────────────────────────────────────
# 2) AI 콘텐츠 생성 (JSON) — 실패해도 {} 반환 (호출부가 템플릿 폴백)
# ─────────────────────────────────────────────────────────────
def _clean_line(s: str) -> str:
    return re.sub(r"\s+", " ", (s or "").replace("*", "").replace("#", "").strip()).strip('"').strip("'")


def _safe_item(s: str, max_len: int = 60) -> str:
    """검증 불가 주장이 섞였으면 버린다(빈 문자열)."""
    s = _clean_line(s)
    if not (2 <= len(s) <= max_len):
        return ""
    low = s.lower()
    if any(b in low for b in _BANNED_CLAIM):
        return ""
    return s


def _clean_list(items, max_len=60, cap=10) -> list:
    out = []
    for it in (items or []):
        v = _safe_item(str(it), max_len)
        if v and v not in out:
            out.append(v)
        if len(out) >= cap:
            break
    return out


def _gemini_report_json(ctx: dict, signals: dict) -> dict:
    if not API_KEY:
        return {}
    if not ctx["site_text"] and not ctx["flagship"]:
        return {}

    market_titles = "\n".join(f"- {t}" for t in signals.get("titles", [])[:24]) or "(검색 데이터 없음)"

    prompt = f"""너는 영유아·임산부 시장 전문 마케팅 기획자다. 형식적인 보고서가 아니라,
너만의 관점·컨셉·아이디어가 살아있는 '기획서'를 쓴다.
아래 '{ctx['brand']}' 브랜드의 실제 제품 정보 + 실제 시장 반응(네이버 글 제목)을 근거로
기획 콘텐츠를 만든다.

--- 홈페이지 텍스트 ---
{ctx['site_text'] or "(없음)"}
--- 주력상품 ---
{ctx['flagship'] or "(없음)"}
--- 주력 상품 라인(핵심어) ---
{signals.get('prod_kw') or "(없음)"}
--- 카테고리 ---
{ctx['category'] or "(없음)"}
--- 실제 시장 반응: 네이버 블로그/카페 글 제목 (브랜드 전반 + 주력 라인, 사람들이 진짜 쓰는 표현) ---
{market_titles}

아래 JSON 형식으로만 답하라(설명·코드펜스 금지):
{{
  "slogan": "이 브랜드를 한 방에 각인시키는 획기적이고 감각적인 슬로건 한 줄 (12~24자)",
  "big_idea": "이 브랜드를 셀픽 채널에서 밀 '캠페인 빅 아이디어' — 컨셉을 관통하는 한 줄 (기획자의 창작)",
  "concept_name": "그 캠페인 컨셉을 부를 짧은 이름 (6~14자, 감각적으로)",
  "insight": "위 실제 글 제목에서 읽어낸 소비자/시장 인사이트 한 줄 (진짜 니즈)",
  "product_summary": "이 제품이 무엇이고 어떤 가치를 주는지 1문장 (홈페이지 근거)",
  "targets": ["이 제품에 해당하는 타겟 세그먼트(아래 목록에서만)"],
  "market_keywords": ["위 실제 글 제목에서 사람들이 실제로 쓰는 표현/니즈 6개 (지어내지 말고 제목에 실제로 드러난 것만)"],
  "appeal_points": ["이 브랜드만의 소구포인트/별칭/바이럴 키워드 5개 (# 없이, 기획자의 창작)"],
  "offline_ideas": ["산부인과/조리원/베이비스튜디오에서의 획기적 오프라인 활용 아이디어 한 줄"],
  "online_ideas": ["맘카페에서 실제로 퍼질 온라인 바이럴 아이디어 한 줄"]
}}

타겟 세그먼트는 반드시 이 목록에서만 고른다: {", ".join(SEGMENT_POOL)}

반드시 지킬 것:
- 홈페이지·주력상품·실제 글 제목에 드러난 것만 근거로. 없는 사실·수치를 지어내지 마라.
- 순위·수상·인증·특허·FDA·제조사명·'1위'·'최고'·구체적 수치 같은 검증 불가 주장 금지
  (자동 폐기됨). 단 slogan·big_idea·concept_name·appeal_points는 '창작 카피'라 자유롭게.
- market_keywords는 반드시 위 '실제 글 제목'에 근거해야 한다 (창작 X).
- big_idea·insight는 뻔하지 않게, 관점이 담기게. offline/online 아이디어 각 3개.
- 3페이지짜리 압축 기획서다. 군더더기 없이 '알맹이'와 '한 방'만."""

    # 간헐적 빈 응답·일시 오류 대비 최대 2회 시도 (실패해도 {} → 템플릿 폴백)
    for _attempt in range(2):
        try:
            r = requests.post(
                API_URL.format(m=MODEL),
                headers={"x-goog-api-key": API_KEY, "Content-Type": "application/json"},
                json={
                    "contents": [{"parts": [{"text": prompt}]}],
                    "generationConfig": {"temperature": 0.6, "maxOutputTokens": 4096},
                },
                timeout=60,
            )
            if r.status_code != 200:
                continue
            parts = r.json()["candidates"][0]["content"]["parts"]
            raw = "".join(p.get("text", "") for p in parts).strip()
        except Exception:
            continue

        # 코드펜스·앞뒤 잡텍스트 제거 후 최외곽 { } 파싱
        raw = re.sub(r"```(?:json)?", "", raw)
        m = re.search(r"\{.*\}", raw, re.DOTALL)
        if not m:
            continue
        try:
            data = json.loads(m.group(0))
        except Exception:
            continue
        if isinstance(data, dict) and data:
            return data
    return {}


# ─────────────────────────────────────────────────────────────
# 3) 콘텐츠 조립 (AI + 안전필터 + 템플릿 폴백)
# ─────────────────────────────────────────────────────────────
# 글 제목 빈출어에서 걸러낼 말 — 표현이 아니라 잡음인 것들
_TITLE_NOISE = {
    "추천", "후기", "리뷰", "내돈내산", "협찬", "체험단", "솔직", "사용", "사용기",
    "정품", "구매", "최저가", "할인", "특가", "이벤트", "증정", "무료", "배송",
    "그리고", "저는", "우리", "너무", "정말", "진짜", "이번", "요즘", "함께",
    "가격", "비교", "정보", "포스팅", "블로그", "카페", "공유", "기록", "일상",
}


# 브랜드 이름에 잘 붙는 조각들. 이걸 품은 토큰은 (상품 유형어가 아닌 한)
# 다른 브랜드명일 확률이 높다 → '시장 표현'에서 제외.
_BRANDISH_SUFFIX = ("베베", "맘마", "키즈", "베이비", "마미", "블리", "뜰", "닷컴", "넷")


# '실제 시장 표현'이라고 말하려면 최소 이만큼의 글은 있어야 한다.
# 그 미만이면 표현이 아니라 몇몇 글의 단편이므로 일반 문구로 대체하고 라벨도 바꾼다.
MIN_TITLES_FOR_REAL_KW = 6


def _keywords_from_titles(titles, brand: str, cap: int = 8, prod_kw: str = "") -> list:
    """AI 없이도 '진짜 글 제목'에서 자주 나오는 표현을 뽑는다 (AI 실패 시 2차 안전망).

    빈도 상위 단어를 그대로 쓰므로 창작이 섞이지 않는다 —
    '실제 후기·글에서 쓰는 표현'이라는 라벨을 지킬 수 있는 최소 장치.

    ⚠ prod_kw를 주면 '그 상품이 실제로 언급된 제목'만 센다.
      네이버는 느슨하게 매칭해서 베이비페어·전시회 글까지 딸려오는데,
      그걸 그대로 세면 경쟁사명(무스텔라)·장소(코엑스)가 '시장 표현'으로 올라간다.
      고객사 제안서에 경쟁사 이름이 실리면 안 된다.
    """
    if not titles:
        return []
    if prod_kw:
        pk = prod_kw.replace(" ", "")
        # 유형어 전체 또는 그 뒷부분(바스앤샴푸 → 샴푸)이 제목에 있으면 그 상품 글로 본다
        tails = {pk} | {pk[i:] for i in range(1, max(1, len(pk) - 1))}
        focused = [t for t in titles
                   if any(x in str(t).replace(" ", "") for x in tails if len(x) >= 2)]
        if len(focused) >= 3:      # 너무 적으면 원래 목록을 쓴다 (표현이 아예 없느니)
            titles = focused
    bn = (brand or "").replace(" ", "")
    counts = {}
    for t in titles:
        for tok in re.split(r"[\s,./|\[\]()\-~!?·]+", str(t)):
            tok = tok.strip()
            # 끝에 붙은 조사 제거 ('바스앤샴푸와' → '바스앤샴푸')
            tok = re.sub(r"(으로|에서|에게|와의|과의|와|과|를|을|은|는|이|가|의|에|도|로)$",
                         "", tok)
            if len(tok) < 2 or len(tok) > 12:
                continue
            if any(ch.isdigit() for ch in tok):
                continue
            if bn and (tok in bn or bn in tok):
                continue
            if tok in _TITLE_NOISE or tok in _PK_STOPWORDS:
                continue
            # 다른 브랜드 이름으로 보이면 버린다 — 고객사 제안서에 경쟁사명이
            # '시장 표현'으로 실리면 안 된다. (몽쉘베베·무스텔라 …)
            #   단 상품 유형어를 품고 있으면 브랜드가 아니라 제품명이다
            #   ('베이비워시'는 살리고 '몽쉘베베'는 버린다)
            if (any(sfx in tok for sfx in _BRANDISH_SUFFIX)
                    and not any(t in tok for t in _PRODUCT_TYPES)):
                continue
            if not re.search(r"[가-힣]", tok):       # 한글이 없는 토큰(영문/기호)은 제외
                continue
            counts[tok] = counts.get(tok, 0) + 1
    # 2회 이상 나온 말만 = 실제로 '반복해서 쓰이는' 표현
    ranked = sorted((w for w, c in counts.items() if c >= 2),
                    key=lambda w: (-counts[w], -len(w)))
    # 조각·중복 정리: '바스앤샴푸'가 이미 있으면 '바스'·'샴푸'·'바스앤샴푸와'는 뺀다
    out = []
    for w in ranked:
        if any(w in kept or kept in w for kept in out):
            continue
        out.append(w)
        if len(out) >= cap:
            break
    return out


def _fallback(ctx: dict) -> dict:
    """AI가 비었을 때 쓰는 카테고리 기반 안전 템플릿 (지어낸 주장 없음)."""
    cat = ctx["category"] or "영유아 제품"
    b = ctx["brand"]
    return {
        "slogan": f"{b}, 엄마의 첫 순간에 함께합니다",
        "big_idea": f"{b}를 '엄마의 첫 선택'으로 각인 — 출산의 시작점에서 브랜드를 심는다",
        "concept_name": "첫 순간 각인",
        "insight": "엄마는 아기에게 '가장 먼저 닿는 것'에 가장 예민하고 지갑을 아끼지 않는다.",
        "product_summary": f"{b}는 영유아·임산부 타깃의 {cat} 브랜드입니다.",
        "targets": ["산모", "아기", "가족"],
        "market_keywords": ["출산준비물", "신생아 필수템", "육아템", "맘카페 추천", "출산가방 리스트"],
        "appeal_points": ["엄마가 고른 브랜드", "믿고 쓰는 데일리템", "선물하기 좋은"],
        "offline_ideas": [
            "산후조리원 신생아실·퇴소 교육 시 실사용·권유",
            "산부인과 대기·상담 공간에서 제품 경험 노출",
            "베이비 스튜디오 방문 고객에게 제품 각인",
        ],
        "online_ideas": [
            "맘스홀릭 등 맘카페 활용법 콘텐츠 바이럴",
            "출산가방 리스트 공유글 필수템 등극",
            "육아 선배맘 후기 기반 추천 확산",
        ],
    }


def build_content(row: dict) -> dict:
    """브랜드 한 행 → 슬라이드에 넣을 콘텐츠 dict (항상 완성형 반환)."""
    ctx = _product_context(row)
    signals = _market_signals(ctx["brand"], ctx["category"], ctx["flagship"])
    ai = _gemini_report_json(ctx, signals)
    fb = _fallback(ctx)

    def _copy(field, max_len):
        """창작 카피(슬로건·빅아이디어 등): 검증불가 주장/과길이만 폐기."""
        v = _clean_line(ai.get(field, "")) if ai else ""
        if not v or len(v) > max_len or any(b in v.lower() for b in _BANNED_CLAIM):
            return fb[field]
        return v

    slogan = _copy("slogan", 34)
    big_idea = _copy("big_idea", 70)
    concept_name = _copy("concept_name", 18)
    insight = _copy("insight", 90)

    # product_summary (안전 필터 후 폴백)
    summary = _clean_line(ai.get("product_summary", "")) if ai else ""
    if not summary or any(b in summary.lower() for b in _BANNED_CLAIM):
        summary = fb["product_summary"]

    # targets — 풀 안에 있는 것만
    targets = [t for t in (ai.get("targets") or []) if t in SEGMENT_POOL]
    if not targets:
        targets = fb["targets"]
    targets = [s for s in SEGMENT_POOL if s in targets][:8]

    # ⭐ '실제 후기·글에서 쓰는 표현'은 진짜 글 제목에서 나와야 한다.
    #   AI가 실패해도 조용히 일반 문구로 바꾸면, 라벨은 '실제'인데 내용은 창작이 된다.
    #   → ① AI(제목 기반) ② 제목에서 직접 빈출어 추출 ③ 그래도 없으면 일반 템플릿
    #   ③일 때만 kw_is_real=False → 페이지가 라벨을 '일반 표현'으로 바꿔 단다.
    #   ⚠ 글이 서너 개뿐이면 '시장 표현'이 아니라 그 몇 글의 단편일 뿐이다.
    #     (누엔쁘: 제목 2개 → '3가지 루틴', '족욕과 경락마사지', '추천 해주세요')
    #     그런 건 '실제 시장 표현'이라고 내밀 수 없다 → 일반 문구 + 라벨 전환.
    _titles = signals.get("titles") or []
    market_kw, kw_is_real = [], False
    if len(_titles) >= MIN_TITLES_FOR_REAL_KW:
        market_kw = _clean_list(ai.get("market_keywords"), 24, 8)
        if not market_kw:
            market_kw = _keywords_from_titles(
                _titles, ctx["brand"], prod_kw=signals.get("prod_kw") or "")
        kw_is_real = bool(market_kw)
    if not market_kw:
        # 일반 문구로 채우지 않는다 — 빈 칸이 되고, 페이지가
        # '아직 노출이 없다'는 문장을 대신 넣는다 (거짓 표현 방지)
        market_kw = []
        kw_is_real = False
    appeal = _clean_list(ai.get("appeal_points"), 24, 8) or fb["appeal_points"]
    offline = _clean_list(ai.get("offline_ideas"), 80, 4) or fb["offline_ideas"]
    online = _clean_list(ai.get("online_ideas"), 80, 5) or fb["online_ideas"]

    return {
        "brand": ctx["brand"],
        "category": ctx["category"],
        "slogan": slogan,
        "big_idea": big_idea,
        "concept_name": concept_name,
        "insight": insight,
        "product_summary": summary,
        "targets": targets,
        "market_keywords": market_kw,
        "market_keywords_real": kw_is_real,   # 진짜 글에서 나왔나 (라벨 문구가 갈린다)
        "appeal_points": appeal,
        "offline_ideas": offline,
        "online_ideas": online,
        "signals": signals,
        "ai_used": bool(ai),
    }


# ─────────────────────────────────────────────────────────────
# 4) PPTX 빌드 — 표지 없이 3페이지, 심플, Pretendard
#    (준 예시 PDF는 '내용/분석' 참고용일 뿐 레이아웃은 새로 잡는다)
# ─────────────────────────────────────────────────────────────
ORANGE = RGBColor(0xF5, 0x82, 0x1F)
ORANGE_D = RGBColor(0xD9, 0x66, 0x0A)
PEACH = RGBColor(0xF3, 0xA9, 0x7A)
PEACH_L = RGBColor(0xFD, 0xF1, 0xE8)
DARK = RGBColor(0x22, 0x22, 0x22)
GRAY = RGBColor(0x66, 0x66, 0x66)
MUTED = RGBColor(0x9A, 0x9A, 0x9A)
LINE = RGBColor(0xE6, 0xE6, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Pretendard"          # 사용자 요청 폰트 (뷰어에 Pretendard 설치 필요)

W = Inches(13.333)
H = Inches(7.5)


def _no_line(shape):
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _rect(slide, l, t, w, h, fill, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    return _no_line(sp)


def _run(p, text, size, bold=False, color=DARK, font=FONT):
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = font
    return r


def _textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def _para(tf, text, size=13, bold=False, color=DARK, align=PP_ALIGN.LEFT,
          first=False, space_after=4, space_before=0, line_spacing=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if line_spacing:
        p.line_spacing = line_spacing
    if text:
        _run(p, text, size, bold, color)
    return p


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _eyebrow(slide, text):
    """페이지 상단 얇은 라벨 + 짧은 오렌지 악센트 바."""
    _rect(slide, Inches(0.72), Inches(0.52), Inches(0.34), Inches(0.055), ORANGE)
    tf = _textbox(slide, Inches(0.72), Inches(0.62), Inches(11.9), Inches(0.32))
    _para(tf, text, 11.5, True, MUTED, first=True)


def _title(slide, text, size=27):
    tf = _textbox(slide, Inches(0.7), Inches(0.92), Inches(11.9), Inches(0.75))
    _para(tf, text, size, True, DARK, first=True)


def _rule(slide, t, l=Inches(0.72), w=Inches(11.9)):
    _rect(slide, l, t, w, Pt(1.0), LINE)


def _note(slide, text, top=Inches(6.95)):
    tf = _textbox(slide, Inches(0.72), top, Inches(11.9), Inches(0.35))
    _para(tf, text, 9.5, False, MUTED, first=True)


def _chip(slide, l, t, text, fill=PEACH, txt=WHITE, size=12):
    w = Inches(0.42 + 0.155 * len(text))
    sp = _rect(slide, l, t, w, Inches(0.44), fill, MSO_SHAPE.ROUNDED_RECTANGLE)
    tfr = sp.text_frame
    tfr.vertical_anchor = MSO_ANCHOR.MIDDLE
    pp = tfr.paragraphs[0]
    pp.alignment = PP_ALIGN.CENTER
    _run(pp, text, size, True, txt)
    return l + w + Inches(0.14)     # 다음 칩 x


def _chips_row(slide, texts, top, x0=Inches(0.72), fill=PEACH, txt=WHITE, size=12):
    """줄바꿈 되는 칩 배치 (폭 넘으면 다음 줄로)."""
    x = x0
    row_t = top
    limit = Inches(12.6)
    for t in texts:
        w = Inches(0.42 + 0.155 * len(t))
        if x + w > limit:
            x = x0
            row_t = row_t + Inches(0.56)
        _chip(slide, x, row_t, t, fill, txt, size)
        x = x + w + Inches(0.14)
    return row_t + Inches(0.56)


def _fmt(n):
    try:
        return f"{int(n):,}"
    except Exception:
        return str(n)


# ── 페이지 1: 컨셉 · 전략 = 슬로건 + 빅 아이디어 + 인사이트 + 타겟 ──
def _page_concept(prs, c):
    s = _blank(prs)
    _eyebrow(s, f"CONCEPT & STRATEGY      셀픽 × {c['brand']}")
    # 히어로 슬로건
    slogan = c["slogan"]
    sz = 40 if len(slogan) <= 16 else (34 if len(slogan) <= 24 else 28)
    tf = _textbox(s, Inches(0.7), Inches(1.08), Inches(11.9), Inches(1.4))
    _para(tf, slogan, sz, True, DARK, first=True, line_spacing=1.05)
    # 브랜드 한 줄 정의
    tf2 = _textbox(s, Inches(0.72), Inches(2.45), Inches(11.9), Inches(0.55))
    _para(tf2, c["product_summary"], 14, False, GRAY, first=True, line_spacing=1.2)
    _rule(s, Inches(3.15))
    # 빅 아이디어 (컨셉)
    tf3 = _textbox(s, Inches(0.72), Inches(3.32), Inches(11.9), Inches(1.05))
    _para(tf3, f"빅 아이디어   「{c['concept_name']}」", 12.5, True, ORANGE_D,
          first=True, space_after=4)
    _para(tf3, c["big_idea"], 16, True, DARK, line_spacing=1.15)
    # 인사이트
    tf4 = _textbox(s, Inches(0.72), Inches(4.75), Inches(11.9), Inches(0.85))
    _para(tf4, "인사이트", 12.5, True, ORANGE_D, first=True, space_after=4)
    _para(tf4, c["insight"], 14, False, GRAY, line_spacing=1.2)
    # 핵심 타겟 칩
    tf5 = _textbox(s, Inches(0.72), Inches(5.85), Inches(11.9), Inches(0.3))
    _para(tf5, "핵심 타겟", 12.5, True, ORANGE_D, first=True)
    _chips_row(s, c["targets"][:6], Inches(6.2))


# ── 페이지 2: 시장 · 니즈 = 실제 검색 데이터 + 키워드 ──
def _page_market(prs, c):
    s = _blank(prs)
    sig = c.get("signals") or {}
    _eyebrow(s, "MARKET & NEEDS")
    _title(s, "시장은 이미 이렇게 말하고 있다")
    _rule(s, Inches(1.72))
    # 실시간 시장 신호 스트립 (진짜 수치) — 브랜드 전반 + 주력 라인
    # 수치는 전부 '브랜드 + 주력 상품 유형어' 기준 (기획서가 그 상품 제안서이므로).
    #   출처는 네이버 블로그·카페 둘뿐 → 앞머리에 한 번 밝힌다.
    parts = []
    kw = sig.get("prod_kw") or ""
    head = f"주력 '{kw}'" if kw else "브랜드"
    if sig.get("blog"):
        parts.append(f"{head} 블로그 {_fmt(sig['blog'])}건")
    if sig.get("cafe"):
        parts.append(f"카페 {_fmt(sig['cafe'])}건")
    if sig.get("blog") and sig.get("cafe"):
        parts.append(f"합계 {_fmt(sig['blog'] + sig['cafe'])}건")
    strip = "   ·   ".join(parts) if parts else "실시간 검색 데이터 없음 (네이버 검색 키 필요)"
    bar = _rect(s, Inches(0.72), Inches(1.95), Inches(11.9), Inches(0.62), PEACH_L,
                MSO_SHAPE.ROUNDED_RECTANGLE)
    btf = bar.text_frame
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    # 문구가 길어지면 한 줄에 안 들어간다 → 글자 크기를 살짝 줄여 넘침 방지.
    #   기준: 폭 11.9in 바에 12.5pt 한글은 대략 100자까지 들어간다(실측 기반).
    _line = f"네이버 블로그·카페 실측   {strip}"
    _sz = 12.5 if len(_line) <= 100 else (11.5 if len(_line) <= 115 else 10.5)
    _run(bp, _line, _sz, True, ORANGE_D)
    # 실제 언급 키워드 (진짜 글 제목 기반)
    tf = _textbox(s, Inches(0.72), Inches(2.95), Inches(11.9), Inches(0.3))
    _kws = [k for k in (c.get("market_keywords") or [])][:8]
    if _kws:
        _para(tf, "실제 후기·글에서 쓰는 표현", 12.5, True, ORANGE_D, first=True)
        nt = _chips_row(s, _kws, Inches(3.32),
                        fill=RGBColor(0xEC, 0xEC, 0xEC), txt=DARK)
    else:
        # 후기가 거의 없는 브랜드 — 일반 문구로 칸을 채우면 거짓이 된다.
        #   비워두고 '노출이 안 되고 있다'는 사실을 그대로 적는다.
        #   이게 곧 셀픽이 필요한 이유이기도 하다.
        _para(tf, "실제 후기·글에서 쓰는 표현", 12.5, True, ORANGE_D, first=True)
        tf0 = _textbox(s, Inches(0.72), Inches(3.32), Inches(11.9), Inches(0.5))
        _para(tf0, "아직 온라인에 쌓인 후기가 거의 없습니다 — "
                   "엄마들이 검색해도 이 제품을 만날 접점이 없는 상태입니다.",
              12, False, DARK, first=True)
        nt = Inches(3.32) + Inches(0.56)
    # 브랜드 소구 키워드 (기획자 창작)
    tf2 = _textbox(s, Inches(0.72), nt + Inches(0.15), Inches(11.9), Inches(0.3))
    _para(tf2, "이 브랜드로 밀 소구 키워드 (제안)", 12.5, True, ORANGE_D, first=True)
    _chips_row(s, [f"#{a}" for a in c["appeal_points"][:5]], nt + Inches(0.55))
    _note(s, "※ 상단 수치는 브랜드명과 주력상품 키워드를 함께 넣어 "
             "네이버 블로그·카페에서 실시간 검색한 실측값이며, 소구 키워드는 셀픽 제안입니다.")


# ── 페이지 3: 셀픽 실행 = 광고 상품 + 실행 아이디어 + 도달 규모 ──
def _page_execution(prs, c):
    s = _blank(prs)
    _eyebrow(s, "SELPIC EXECUTION")
    _title(s, "셀픽으로 이렇게 실행한다")
    _rule(s, Inches(1.72))
    # 좌: 셀픽 광고 상품 5종
    tfL = _textbox(s, Inches(0.72), Inches(1.95), Inches(5.7), Inches(4.4))
    _para(tfL, "셀픽 광고 상품", 12.5, True, ORANGE_D, first=True, space_after=7)
    for name, desc in SELPIC_MEDIA:
        _para(tfL, f"● {name}", 12.5, True, DARK, space_after=1, space_before=3)
        _para(tfL, f"    {desc}", 11, False, GRAY, space_after=2)
    # 우: 실행 아이디어 (오프+온)
    tfR = _textbox(s, Inches(6.8), Inches(1.95), Inches(5.9), Inches(4.4))
    _para(tfR, "오프라인 실행", 12.5, True, ORANGE_D, first=True, space_after=5)
    for idea in c["offline_ideas"][:3]:
        _para(tfR, f"● {idea}", 11.5, False, DARK, space_after=4)
    _para(tfR, "온라인 실행", 12.5, True, ORANGE_D, space_before=8, space_after=5)
    for idea in c["online_ideas"][:3]:
        _para(tfR, f"● {idea}", 11.5, False, DARK, space_after=4)
    # 셀픽 도달 규모 (고정 수치 스트립)
    bar = _rect(s, Inches(0.72), Inches(6.05), Inches(11.9), Inches(0.66), PEACH_L,
                MSO_SHAPE.ROUNDED_RECTANGLE)
    btf = bar.text_frame
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    btf.word_wrap = True
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    _run(bp, "연 10만 신생아·부모 DB   ·   전국 460대 키오스크   ·   "
             "조리원 230곳(50% M/S)   ·   조리원 체류 14일 매일 접촉",
         11.5, True, ORANGE_D)
    _note(s, "※ 실행 아이디어는 셀픽 제안이며, 도달 규모는 셀픽 미디어 실측치입니다. 물량은 협의 후 확정.")


def build_pptx(content: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    _page_concept(prs, content)
    _page_market(prs, content)
    _page_execution(prs, content)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─────────────────────────────────────────────────────────────
# 5) 오케스트레이션
# ─────────────────────────────────────────────────────────────
def _safe_filename(brand: str) -> str:
    base = re.sub(r'[\\/:*?"<>|]', "", (brand or "브랜드")).strip() or "브랜드"
    return f"{base}_셀픽MediaProposal.pptx"


def make_report(row: dict):
    """브랜드 행(dict) → (pptx_bytes, filename). AI 실패해도 템플릿으로 완성."""
    content = build_content(row)
    data = build_pptx(content)
    return data, _safe_filename(content["brand"])
